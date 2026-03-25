import re
import qrcode
from pathlib import Path
from typing import Dict, Any, Optional
import json
import numpy as np

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE, MSO_VERTICAL_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn
from PIL import Image
import pdb
from src.state.poster_state import PosterState
from utils.src.logging_utils import log_agent_info, log_agent_success, log_agent_error
from src.config.poster_config import load_config

class Renderer:
    """ 渲染器类，用于设置结构布局，同时根据布局渲染海报 """

    def __init__(self):
        self.name = "renderer"
        # load configuration
        self.config = load_config()
        self.indentation_config = self.config["indentation"]

    def __call__(self, state: PosterState) -> PosterState:
        """ 一个叫posterlayout,一个叫sectionlayout """
        # load sections information
        section_info = state["story_board"]#这里应该时story_board的数据
        section_info = self._preprocess_section_info(state, section_info)#预处理，将section_content改为提取内容
        # load poster layout data
        print("start_select_poster_layout")
        with open(Path(state["resource_dir"]) / "poster_layouts/new_poster_layouts.json", 'r', encoding='utf-8') as f:
            poster_layouts = json.load(f)
            poster_layouts = poster_layouts[str(state["section_number"])]
        # 选取合适的poster_layout
        for poster_layout in poster_layouts:
            print(f"start_set_layout{poster_layout['id']}")
            sections_layout = self._set_layout(state, section_info, poster_layout["layout"])
            sections_layout = self._add_poster_margin_info(sections_layout)
            color_scheme = state["color_scheme"]
            elements_layout = self._set_elements_layout(sections_layout, color_scheme)
            self._save_elements_info(state, poster_layout['id'],elements_layout)
            # render poster
            print("start_render_poster")
            self.render_poster(state, poster_layout['id'], elements_layout)
        self._打分()

        return state
    
    def _save_elements_info(self, state: PosterState, poster_layout_id: int, elements_layout: list):
        output_dir = Path(state["output_dir"])
        with open(output_dir/"content"/f"elements_layout_{poster_layout_id}.json", "w", encoding="utf-8") as f:
            json.dump(elements_layout, f, indent=2)

    def _preprocess_section_info(self, state:PosterState, section_info: Dict):
        new_section_info = section_info.copy()
        sections = new_section_info["spatial_content_plan"]["sections"]
        title_author = state["narrative_content"]["meta"]
        for section in sections:
            if section["section_title"] == "title_author":
                section["text_content"] = [title_author["poster_title"], title_author["authors"]]
        state["story_board"] = new_section_info
        return new_section_info

    def _set_layout(self, state: PosterState, section_info: Dict, poster_layout: list) -> Dict:
        ratio = 3 #海报英寸单位比虚拟16：9单位大3倍
        sections_layout = []
        sections = section_info["spatial_content_plan"]["sections"]
        for section, section_location in zip(sections, poster_layout):
            section_layout = {
                "section_title": section["section_title"],
                "section_content": section["text_content"],
                "x": section_location[0]*ratio,
                "y": section_location[1]*ratio,
                "width": (section_location[2] - section_location[0])*ratio,
                "height": (section_location[3] - section_location[1])*ratio
            }
            visuals_layout = []
            for item in section["visual_assets"]:
                visual_id = item["visual_id"]
                visual_path, aspect = self._get_visual_path(visual_id, state)
                if visual_path and Path(visual_path).exists():
                    visual_layout = {
                        "path":visual_path,
                        "x":section_location[0],
                        "y":section_location[1],
                        "width":3,#后续需要一定的变化
                        "height":3/aspect,
                        "aspect":aspect,
                    }
                    visuals_layout.append(visual_layout)
            section_layout["visuals_layout"] = visuals_layout
            sections_layout.append(section_layout)
        return sections_layout

    def _get_visual_path(self, visual_id: str, state: PosterState):
        """get path to visual asset"""
        images = state.get("images", {})
        tables = state.get("tables", {})
        vid = (visual_id or "").split('_')[-1]
        
        if visual_id.startswith("figure"):
            return images.get(vid, {}).get("path"), images.get(vid, {}).get("aspect")
        if visual_id.startswith("table"):
            return tables.get(vid, {}).get("path"), tables.get(vid, {}).get("aspect")
        
        return None, None

    def _add_poster_margin_info(self, sections_layout: list) -> list:
        poster_margin = self.config["layout"]["poster_margin"]
        for section_layout in sections_layout:
            section_layout["x"] += poster_margin
            section_layout["y"] += poster_margin
        return sections_layout
    
    def _set_elements_layout(self, sections_layout: list, color_scheme: Dict) -> list:
        elements_layout = []
        for section_layout in sections_layout:
            if section_layout["section_title"] == "title_author":
                text_info = {
                    "title_font_size": self.config["typography"]["sizes"]["title"],
                    "author_font_size": self.config["typography"]["sizes"]["authors"],
                    "title_font": self.config["typography"]["fonts"]["title"],
                    "author_font": self.config["typography"]["fonts"]["authors"],
                    "title_font_color":RGBColor(0, 0, 0),
                    "author_font_color":RGBColor(60, 60, 60),
                }
                element_layout = {
                    "element_type": "title_author",
                    "content": section_layout["section_content"],
                    "x": section_layout["x"],
                    "y": section_layout["y"],
                    "width": section_layout["width"],
                    "height": section_layout["height"],
                    "border_width": 0.5,
                    "border_color": self._parse_color(color_scheme["theme"]),
                    **text_info
                }
                elements_layout.append(element_layout)
            else:
                x = section_layout["x"]
                y = section_layout["y"]
                width = section_layout["width"]
                height = section_layout["height"]
                # 处理重点板块背景色
                section_title = section_layout["section_title"]
                subtitle_item = section_title.split("_")
                if "core" in subtitle_item or "results" in subtitle_item or "findings" in subtitle_item:
                    section_container_element = {
                        "element_type": "section_container",
                        "x": x + self.config["layout"]["section_margin"],
                        "y": y + self.config["layout"]["section_margin"],
                        "width": width - self.config["layout"]["section_margin"] * 2,
                        "height": height - self.config["layout"]["section_margin"] * 2,
                        "color": self._parse_color(color_scheme["mono_light"])
                    }
                    elements_layout.append(section_container_element)
                # 处理子标题
                if len(subtitle_item) > 2:
                    subtitle_content = "_".join(subtitle_item[:2])
                else:
                    subtitle_content = "_".join(subtitle_item)
                subtitle_element_layout = {
                    "element_type": "subtitle",
                    "section_title": section_title,
                    "content": subtitle_content,
                    "x": x + self.config["layout"]["subtitle_margin"],
                    "y": y + self.config["layout"]["subtitle_margin"],
                    "width": width - self.config["layout"]["subtitle_margin"] * 2,
                    "height": self.config["typography"]["sizes"]["section_title"] / 72 + self.config["layout"]["subtitle_margin"],
                    "font_size": self.config["typography"]["sizes"]["section_title"],
                    "font": self.config["typography"]["fonts"]["section_title"]
                }
                elements_layout.append(subtitle_element_layout)
                y += self.config["typography"]["sizes"]["section_title"] / 72 + self.config["layout"]["subtitle_margin"] * 2
                height -= self.config["typography"]["sizes"]["section_title"] / 72 + self.config["layout"]["subtitle_margin"] * 2

                # 处理图片
                visual_assets_num = len(section_layout["visuals_layout"])
                visual_margin = self.config["layout"]["visual"]["margin"]
                visual_spacing = self.config["layout"]["visual"]["spacing"]
                visual_max_width = self.config["layout"]["visual"]["max_width"]
                if visual_assets_num == 0:
                    pass
                elif visual_assets_num == 1:
                    aspect = section_layout["visuals_layout"][0]["aspect"]
                    if width / height > 1.8:
                        visual_height = min(height - 2*visual_margin, visual_max_width / aspect)
                        visual_element_layout = {
                            "element_type": "visual",
                            "path": section_layout["visuals_layout"][0]["path"],
                            "x": x + visual_margin,
                            "y": y + (height - visual_height) / 2 ,
                            "width": visual_height * aspect,
                            "height": visual_height
                        }
                        x += visual_height * aspect + visual_margin*2
                        width -= visual_height * aspect + visual_margin*2
                        elements_layout.append(visual_element_layout)
                    else:
                        visual_width = min(width - 2*visual_margin, visual_max_width)
                        visual_element_layout = {
                            "element_type": "visual",
                            "path": section_layout["visuals_layout"][0]["path"],
                            "x": x + (width - visual_width) / 2 ,
                            "y": y + visual_margin,
                            "width": visual_width,
                            "height": visual_width / aspect
                        }
                        y += visual_width / aspect + visual_margin*2
                        height -= visual_width / aspect + visual_margin*2
                        elements_layout.append(visual_element_layout)
                elif visual_assets_num == 2:
                    visual_asset1 = section_layout["visuals_layout"][0]
                    visual_path1 = visual_asset1["path"]
                    visual_aspect1 = visual_asset1["aspect"]
                    visual_asset2 = section_layout["visuals_layout"][1]
                    visual_path2 = visual_asset2["path"]
                    visual_aspect2 = visual_asset2["aspect"]
                    max_text_space = 0
                    best_strategy = None
                    reciprocal_aspect_sum = 1 / visual_aspect1 + 1 / visual_aspect2
                    #两张图片靠左摆放上下叠放
                    if (visual_max_width / visual_aspect1 + visual_max_width / visual_aspect2) * 0.7 + visual_spacing + visual_margin*2 <= height:
                        visual_width = min((height - 2*visual_margin - visual_spacing) / reciprocal_aspect_sum, visual_max_width)
                        text_space = (width - visual_width) * height
                        if text_space > max_text_space:
                            max_text_space = text_space
                            best_strategy = "left1"
                    # 两张图片靠左摆放左右摆放
                    if (visual_max_width * 2 + visual_spacing + visual_margin*2) * 0.7 <= width:
                        visual_width = min((height - 2*visual_margin) * min(visual_aspect1, visual_aspect2)*2, visual_max_width*2)
                        text_space = (width - visual_width) * height
                        if text_space > max_text_space:
                            max_text_space = text_space
                            best_strategy = "left2"
                    # 两张照片靠上摆放左右摆放
                    if (visual_max_width * 2 + visual_spacing + visual_margin*2) * 0.7 <= width:
                        visual_height = min((width - 2*visual_margin - visual_spacing)/2, visual_max_width) / min(visual_aspect1, visual_aspect2)
                        text_space = (height - visual_height) * width
                        if text_space > max_text_space:
                            max_text_space = text_space
                            best_strategy = "top1"
                    # 两张照片靠上摆放上下摆放
                    if (visual_max_width / visual_aspect1 + visual_max_width / visual_aspect2) * 0.7 + visual_spacing + visual_margin*2 <= height:
                        visual_height = min(width - 2*visual_margin, visual_max_width) * reciprocal_aspect_sum
                        text_space = (height - visual_height) * width
                        if text_space > max_text_space:
                            max_text_space = text_space
                            best_strategy = "top2"
                    if best_strategy == "left1":
                        visual_width = min((height - 2*visual_margin - visual_spacing) / reciprocal_aspect_sum, visual_max_width)
                        visual_total_height = visual_width * reciprocal_aspect_sum + visual_spacing
                        visual_element_layout1 = {
                            "element_type": "visual",
                            "path": visual_path1,
                            "x": x + visual_margin,
                            "y": y + (height - visual_total_height) / 2 ,
                            "width": visual_width,
                            "height": visual_width / visual_aspect1
                        }
                        visual_element_layout2 = {
                            "element_type": "visual",
                            "path": visual_path2,
                            "x": x + visual_margin + visual_width,
                            "y": visual_element_layout1["y"] + visual_element_layout1["height"] + visual_spacing,
                            "width": visual_width,
                            "height": visual_width / visual_aspect2
                        }
                        elements_layout.append(visual_element_layout1)
                        elements_layout.append(visual_element_layout2)
                        x += (visual_width + 2*visual_margin)
                        width -= (visual_width + 2*visual_margin)
                    elif best_strategy == "left2":
                        visual_width = min((height - 2*visual_margin) * min(visual_aspect1, visual_aspect2)*2, visual_max_width*2)
                        visual_height1 = visual_width / visual_aspect1
                        visual_element_layout1 = {
                            "element_type": "visual",
                            "path": visual_path1,
                            "x": x + visual_margin,
                            "y": y + (height - visual_height1) / 2 ,
                            "width": visual_width,
                            "height": visual_height1
                        }
                        visual_height2 = visual_width / visual_aspect2
                        visual_element_layout2 = {
                            "element_type": "visual",
                            "path": visual_path2,
                            "x" : visual_element_layout1["x"] + visual_element_layout1["width"] + visual_spacing,
                            "y": y + (height - visual_height2) / 2 ,
                            "width": visual_width,
                            "height": visual_height2
                        }
                        elements_layout.append(visual_element_layout1)
                        elements_layout.append(visual_element_layout2)
                        x += (visual_width*2 + visual_margin*2 + visual_spacing)
                        width -= (visual_width*2 + visual_margin*2 + visual_spacing)
                    elif best_strategy == "top1":
                        visual_width = min((width - 2*visual_margin - visual_spacing)/2, visual_max_width)
                        visual_height = visual_width / min(visual_aspect1, visual_aspect2)
                        visual_total_width = visual_width*2 + visual_spacing
                        visual_element_layout1 = {
                            "element_type": "visual",
                            "path": visual_path1,
                            "x": x + (width - visual_total_width) / 2 ,
                            "y": y + visual_margin,
                            "width": visual_width,
                            "height": visual_width / visual_aspect1
                        }
                        visual_element_layout2 = {
                            "element_type": "visual",
                            "path": visual_path2,
                            "x": visual_element_layout1["x"] + visual_element_layout1["width"] + visual_spacing,
                            "y": y + visual_margin,
                            "width": visual_width,
                            "height": visual_width / visual_aspect2
                        }
                        elements_layout.append(visual_element_layout1)
                        elements_layout.append(visual_element_layout2)
                        y += (visual_height + 2*visual_margin)
                        height -= (visual_height + 2*visual_margin)
                    elif best_strategy == "top2":
                        visual_width = min(width - 2*visual_margin, visual_max_width)
                        visual_element_layout1 = {
                            "element_type": "visual",
                            "path": visual_path1,
                            "x": x + (width - visual_width) / 2 ,
                            "y": y + visual_margin,
                            "width": visual_width,
                            "height": visual_width / visual_aspect1
                        }
                        visual_element_layout2 = {
                            "element_type": "visual",
                            "path": visual_path2,
                            "x": x + (width - visual_width) / 2 ,
                            "y": visual_element_layout1["y"] + visual_element_layout1["height"] + visual_spacing,
                            "width": visual_width,
                            "height": visual_width / visual_aspect2
                        }
                        elements_layout.append(visual_element_layout1)
                        elements_layout.append(visual_element_layout2)
                        y += (visual_element_layout1["height"] + visual_element_layout2["height"] + visual_margin*2 + visual_spacing)
                        height -= (visual_element_layout1["height"] + visual_element_layout2["height"] + visual_margin*2 + visual_spacing)
                # 处理文本
                text_margin = self.config["layout"]["text_margin"]
                text_element_layout = {
                    "element_type": "text",
                    "content": section_layout["section_content"],
                    "x": x + text_margin,
                    "y": y + text_margin,
                    "width": width - text_margin*2,
                    "height": height - text_margin*2,
                    "font_size": self.config["typography"]["sizes"]["body_text"],
                    "font": self.config["typography"]["fonts"]["body_text"]
                }
                elements_layout.append(text_element_layout)
        return elements_layout

    def render_poster(self, state: PosterState,poster_layout_id: int , elements_layout: list):
        """ 渲染海报 """
        prs = Presentation()
        prs.slide_width = Inches(state["poster_width"])
        prs.slide_height = Inches(state["poster_height"])
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        for element_layout in elements_layout:
            if element_layout["element_type"] == "title_author":
                textbox = slide.shapes.add_textbox(
                    Inches(element_layout["x"]),
                    Inches(element_layout["y"]),
                    Inches(element_layout["width"]),
                    Inches(element_layout["height"])
                )
                line = textbox.line
                line.fill.solid()
                line.fill.fore_color.rgb = element_layout["border_color"]
                line.width = Pt(element_layout["border_width"])
                tf = textbox.text_frame
                tf.clear()  # 清空默认文本
                tf.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE  # 垂直锚点（text_frame属性）
                tf.word_wrap = True  # 自动换行（新增，避免文本溢出）
                if element_layout["height"] / element_layout["width"] > 2:
                    bodyPr = tf._txBody.bodyPr
                    bodyPr.set(qn('a:vert'), 'vert270')
                title_para = tf.add_paragraph()  # 重新创建标题段落（避免默认段落问题）
                title_para.text = element_layout["content"][0]
                title_para.font.size = Pt(element_layout["title_font_size"])
                title_para.font.bold = True
                title_para.alignment = PP_ALIGN.CENTER
                title_para.font.name = element_layout["title_font"]
                title_para.font.color.rgb = element_layout["title_font_color"]
                author_para = tf.add_paragraph()  # 重新创建作者段落（避免默认段落问题）
                author_para.text = element_layout["content"][1]
                author_para.font.size = Pt(element_layout["author_font_size"])
                author_para.alignment = PP_ALIGN.CENTER
                author_para.font.name = element_layout["author_font"]
                author_para.font.color.rgb = element_layout["author_font_color"]
            elif element_layout["element_type"] == "section_container":
                container = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(element_layout["x"]), Inches(element_layout["y"]), Inches(element_layout["width"]), Inches(element_layout["height"]))
                container.fill.solid()
                container.fill.fore_color.rgb = element_layout["color"]
            elif element_layout["element_type"] == "subtitle":
                textbox = slide.shapes.add_textbox(
                    Inches(element_layout["x"]),
                    Inches(element_layout["y"]),
                    Inches(element_layout["width"]),
                    Inches(element_layout["height"])
                )
                tf = textbox.text_frame
                tf.clear()  # 清空默认文本
                tf.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP  # 垂直锚点（text_frame属性）
                tf.word_wrap = False  # 不自动换行（新增，避免文本溢出）
                subtitle_para = tf.add_paragraph()  # 重新创建子标题段落（避免默认段落问题）
                subtitle_para.text = element_layout["content"]
                subtitle_para.font.size = Pt(element_layout["font_size"])
                subtitle_para.alignment = PP_ALIGN.LEFT
                subtitle_para.font.name = element_layout["font"]
                subtitle_para.font.bold = True
            elif element_layout["element_type"] == "visual":
                image = slide.shapes.add_picture(
                    element_layout["path"],
                    Inches(element_layout["x"]),
                    Inches(element_layout["y"]),
                    Inches(element_layout["width"]),
                    Inches(element_layout["height"])
                )
            elif element_layout["element_type"] == "text":
                textbox = slide.shapes.add_textbox(
                    Inches(element_layout["x"]),
                    Inches(element_layout["y"]),
                    Inches(element_layout["width"]),
                    Inches(element_layout["height"])
                )
                textbox.fill.solid()
                textbox.fill.fore_color.rgb = RGBColor(255, 255, 255)
                textbox.fill.fore_color.transparency = 1.0  # 完全透明
                textbox.line.fill.background()  # 无边框

                tf = textbox.text_frame
                tf.clear()  # 清空默认文本
                tf.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP  # 垂直锚点（text_frame属性）
                tf.word_wrap = True
                for line in element_layout["content"]:
                    content_para = tf.add_paragraph()
                    line = line.strip()
                    if line.strip().startswith(self.indentation_config["secondary_bullet_char"]):
                        content_para.level = self.indentation_config["secondary_level"]
                    else:
                        content_para.level = self.indentation_config["primary_level"]
                    self._add_formatted_runs(content_para, line, element_layout["font"], Pt(element_layout["font_size"]), RGBColor(0, 0, 0))
                    content_para.alignment = PP_ALIGN.LEFT
        prs.save(Path(state["output_dir"]) / f"poster_{poster_layout_id}.pptx")
        print(f"save poster{poster_layout_id}")
        return 0
    
    def _add_formatted_runs(self, paragraph, text: str, font_family: str, 
                          base_font_size, base_color):
        """parse text and create separate runs for each format type"""
        # tokenize the text into formatting segments
        segments = self._tokenize_formatting(text)
        
        # create runs for each segment
        for segment in segments:
            run = paragraph.add_run()
            run.text = segment['text']
            run.font.name = font_family
            run.font.size = base_font_size
            
            # apply formatting based on segment type
            if segment['color']:
                run.font.color.rgb = self._parse_color(segment['color'])
            else:
                run.font.color.rgb = base_color
            
            if segment['bold']:
                run.font.bold = True
            
            if segment['italic']:
                run.font.italic = True

    def _tokenize_formatting(self, text: str) -> list:
        """tokenize text into formatting segments with precise position tracking"""
        segments = []
        i = 0
        
        while i < len(text):
            # check for color markup: <color:#RRGGBB>text</color>
            color_match = re.match(r'<color:(#[0-9A-Fa-f]{6})>', text[i:])
            if color_match:
                color_hex = color_match.group(1)
                opening_tag_end = i + color_match.end()
                
                # find closing </color> tag using absolute position
                closing_tag_pattern = r'</color>'
                color_content_start = opening_tag_end
                closing_match = re.search(closing_tag_pattern, text[color_content_start:])
                
                if closing_match:
                    # calculate absolute positions
                    color_content_end = color_content_start + closing_match.start()
                    closing_tag_end = color_content_start + closing_match.end()
                    
                    # extract content between color tags
                    colored_text = text[color_content_start:color_content_end]
                    
                    # process colored text with automatic bold
                    if colored_text.strip():  # only process non-empty content
                        segments.append({
                            'text': colored_text,
                            'bold': True,  # all colored text is bold
                            'italic': False,
                            'color': color_hex
                        })
                    
                    # move past the entire color block
                    i = closing_tag_end
                    continue
                else:
                    # malformed color tag, treat as regular text
                    segments.append({
                        'text': text[i],
                        'bold': False,
                        'italic': False,
                        'color': None
                    })
                    i += 1
                    continue
            
            # check for bold: **text**
            bold_match = re.match(r'\*\*(.*?)\*\*', text[i:])
            if bold_match:
                bold_text = bold_match.group(1)
                segments.append({
                    'text': bold_text,
                    'bold': True,
                    'italic': False,
                    'color': None
                })
                i += bold_match.end()
                continue
            
            # check for italic: *text*
            italic_match = re.match(r'\*(.*?)\*', text[i:])
            if italic_match:
                italic_text = italic_match.group(1)
                segments.append({
                    'text': italic_text,
                    'bold': False,
                    'italic': True,
                    'color': None
                })
                i += italic_match.end()
                continue
            
            # regular text - find next formatting marker
            next_format = re.search(r'(\*\*|\*|<color:)', text[i:])
            if next_format:
                regular_text = text[i:i + next_format.start()]
            else:
                regular_text = text[i:]
            
            if regular_text:
                segments.append({
                    'text': regular_text,
                    'bold': False,
                    'italic': False,
                    'color': None
                })
            
            if next_format:
                i += next_format.start()
            else:
                break
        
        return segments

    def _parse_color(self, color_str: str) -> RGBColor:
        """parse color string to RGBColor"""
        hex_color = color_str.lstrip('#')
        r, g, b = (int(hex_color[i:i+2], 16) for i in (0, 2, 4))
        return RGBColor(r, g, b)

def renderer_node(state: PosterState) -> Dict[str, Any]:
    print("start_render_node")
    result = Renderer()(state)
    return {
        **state,
        "tokens": result["tokens"],
        "current_agent": result["current_agent"],
        "errors": result["errors"]
    }